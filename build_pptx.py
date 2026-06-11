"""Generate the AI Agent Program roadmap PowerPoint deck (AWS, 5 months)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- Palette ----------
AWS_NAVY   = RGBColor(0x23, 0x2F, 0x3E)   # dark slate
AWS_ORANGE = RGBColor(0xFF, 0x99, 0x00)   # AWS orange
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT      = RGBColor(0xF2, 0xF4, 0xF7)
GREY       = RGBColor(0x6B, 0x74, 0x80)
DARK       = RGBColor(0x1B, 0x24, 0x30)

PHASE_COLORS = [
    RGBColor(0x2E, 0x86, 0xC1),  # Phase 1 - blue
    RGBColor(0x28, 0xB4, 0x63),  # Phase 2 - green
    RGBColor(0x8E, 0x44, 0xAD),  # Phase 3 - purple
    RGBColor(0xE6, 0x7E, 0x22),  # Phase 4 - orange
    RGBColor(0x16, 0xA0, 0x85),  # Phase 5 - teal
    RGBColor(0xC0, 0x39, 0x2B),  # Phase 6 - red
]

# (name, objective, start_month(1-idx), duration_months, aws_services)
# Compressed 5-month plan - phases run in parallel to hit the deadline.
PHASES = [
    ("Phase 1 - Administrative AI Agent (MVP)",
     "Automate meeting admin: transcripts to minutes, actions, owners, due dates",
     1, 1,
     "Amazon Transcribe | Amazon Bedrock | Amazon S3 | DynamoDB | AWS Lambda"),
    ("Phase 2 - Action Tracking Agent",
     "Monitor, remind, escalate and report on open / overdue actions",
     2, 1,
     "Amazon EventBridge | Amazon SES / SNS | DynamoDB | AWS Lambda"),
    ("Phase 3 - Governance & Forum Agent",
     "Read SteerCo / CIO / Cyber / Risk / Audit outputs; detect themes & risks",
     2, 2,
     "Amazon Bedrock | Bedrock Knowledge Bases | Amazon S3 | AWS Lambda"),
    ("Phase 4 - Reporting & Analytics Agent",
     "Consolidate reports, measure delivery, build CIO packs & trend analysis",
     3, 2,
     "Amazon QuickSight | AWS Glue | Amazon Athena | Amazon S3"),
    ("Phase 5 - Enterprise Knowledge Agent",
     "Searchable organisational memory; natural-language Q&A over all outputs",
     4, 1,
     "Bedrock Knowledge Bases | OpenSearch Serverless | Amazon Bedrock (RAG)"),
    ("Phase 6 - Central Control Tower (Orchestration)",
     'Connect all agents through a central platform - the "spine"',
     4, 2,
     "Amazon Bedrock Agents | AWS Step Functions | API Gateway | Amazon Cognito"),
]

MONTHS = [f"M{i}" for i in range(1, 6)]
N_MONTHS = len(MONTHS)

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, l, t, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def add_text(slide, l, t, w, h, text, size, color, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = align
    for i, line in enumerate(text.split("\n")):
        run_p = p if i == 0 else tf.add_paragraph()
        run_p.alignment = align
        r = run_p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    return tb


# ============================================================
# SLIDE 1 - Title
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, AWS_NAVY)
add_rect(s, 0, Inches(5.05), SW, Inches(0.10), AWS_ORANGE)
add_text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.4),
         "Enterprise AI Agent Programme", 44, WHITE, bold=True)
add_text(s, Inches(0.9), Inches(3.2), Inches(11.5), Inches(0.9),
         "Delivery Roadmap: MVP to Central Control Tower", 26, AWS_ORANGE, bold=True)
add_text(s, Inches(0.92), Inches(5.25), Inches(11.5), Inches(0.6),
         "6 Phases  -  5 Months  -  Built on Amazon Web Services (AWS)", 18, LIGHT)
add_text(s, Inches(0.92), Inches(6.7), Inches(11.5), Inches(0.4),
         "Powered by Amazon Bedrock, Transcribe, QuickSight & Step Functions",
         12, GREY, italic=True)

# ============================================================
# SLIDE 2 - Gantt roadmap
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
add_rect(s, 0, 0, SW, Inches(0.85), AWS_NAVY)
add_text(s, Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.6),
         "Programme Roadmap - 5 Month Timeline", 26, WHITE, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)

# layout geometry
label_l = Inches(0.4)
label_w = Inches(3.9)
grid_l  = Inches(4.4)
grid_r  = Inches(13.0)
grid_w  = grid_r - grid_l
top     = Inches(1.55)
col_w   = Emu(int(grid_w / N_MONTHS))
row_h   = Inches(0.72)
row_gap = Inches(0.08)

# month header
add_text(s, label_l, Inches(1.0), label_w, Inches(0.45), "Phase / Workstream",
         12, DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
for i, m in enumerate(MONTHS):
    cl = Emu(int(grid_l) + i * int(col_w))
    add_rect(s, cl, Inches(1.0), col_w, Inches(0.45), LIGHT)
    add_text(s, cl, Inches(1.0), col_w, Inches(0.45), m, 11, DARK, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# vertical gridlines + quarter bands behind rows
total_h = Emu(int(row_h) * 6 + int(row_gap) * 5)
for i in range(N_MONTHS + 1):
    cl = Emu(int(grid_l) + i * int(col_w))
    add_rect(s, cl, top, Emu(9525), total_h, RGBColor(0xE2, 0xE6, 0xEA))

for idx, (name, obj, start, dur, aws) in enumerate(PHASES):
    rt = Emu(int(top) + idx * (int(row_h) + int(row_gap)))
    # row label
    add_text(s, label_l, rt, label_w, row_h, name, 11, DARK, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    # bar
    bl = Emu(int(grid_l) + (start - 1) * int(col_w))
    bw = Emu(dur * int(col_w))
    bar = add_rect(s, Emu(int(bl) + 9525), Emu(int(rt) + 95250),
                   Emu(int(bw) - 19050), Emu(int(row_h) - 190500),
                   PHASE_COLORS[idx], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    span = f"M{start}" if dur == 1 else f"M{start}-M{start + dur - 1}"
    add_text(s, bl, rt, bw, row_h, span, 10, WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# milestone callouts under grid
add_rect(s, grid_l, Inches(6.35), grid_w, Inches(0.02), AWS_ORANGE)
milestones = [
    ("M1", "MVP live"),
    ("M3", "Governance insights"),
    ("M4", "Knowledge Q&A"),
    ("M5", "Control Tower"),
]
for m, lbl in milestones:
    i = int(m[1:]) - 1
    cl = Emu(int(grid_l) + i * int(col_w))
    dia = Inches(0.16)
    add_rect(s, Emu(int(cl) + int(col_w) - int(dia)//2 - 9525), Inches(6.27),
             dia, dia, AWS_ORANGE, shape=MSO_SHAPE.OVAL)
    add_text(s, Emu(int(cl) - int(col_w)), Inches(6.5), Emu(int(col_w)*2),
             Inches(0.5), f"{m}: {lbl}", 9, DARK, bold=True, align=PP_ALIGN.CENTER)

# legend note
add_text(s, label_l, Inches(6.95), Inches(12.5), Inches(0.4),
         "Diamonds = key milestones.  Bars span planned delivery months per phase.",
         10, GREY, italic=True)

# ============================================================
# SLIDE 3 - Phase detail + AWS services
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
add_rect(s, 0, 0, SW, Inches(0.85), AWS_NAVY)
add_text(s, Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.6),
         "Phase Detail & AWS Service Mapping", 26, WHITE, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)

col_x = [Inches(0.4), Inches(4.55), Inches(9.7)]
hdrs = ["Phase & Timeline", "Objective", "Core AWS Services"]
col_w3 = [Inches(4.0), Inches(5.0), Inches(3.2)]
htop = Inches(1.05)
for x, w, h in zip(col_x, col_w3, hdrs):
    add_rect(s, x, htop, w, Inches(0.45), AWS_ORANGE)
    add_text(s, x, htop, w, Inches(0.45), h, 12, WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

rtop = Inches(1.6)
rh = Inches(0.88)
for idx, (name, obj, start, dur, aws) in enumerate(PHASES):
    y = Emu(int(rtop) + idx * int(rh))
    band = LIGHT if idx % 2 == 0 else WHITE
    add_rect(s, col_x[0], y, Inches(12.6), rh, band)
    # colour chip
    add_rect(s, col_x[0], y, Inches(0.12), rh, PHASE_COLORS[idx])
    span = f"M{start}" if dur == 1 else f"M{start}-M{start + dur - 1}"
    add_text(s, Emu(int(col_x[0]) + 120000), y, Inches(3.85), rh,
             f"{name}\n({span})", 10.5, DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, col_x[1], y, col_w3[1], rh, obj, 10, DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, col_x[2], y, col_w3[2], rh, aws.replace(" | ", "\n"), 9, GREY,
             anchor=MSO_ANCHOR.MIDDLE)

add_text(s, col_x[0], Inches(6.95), Inches(12.6), Inches(0.4),
         "Cross-cutting from M1: Security & IAM, Cost governance, Data privacy, CI/CD, Human-in-the-loop review.",
         10, GREY, italic=True)

prs.save(os.path.join(os.path.dirname(os.path.abspath(__file__)), "AI_Agent_Programme_Roadmap.pptx"))
print("PPTX saved:", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
